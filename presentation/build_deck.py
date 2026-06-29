"""
Generates the brownbag presentation:
    "Identifying Backend Bottlenecks — Technical & Process"

Run:
    python presentation/build_deck.py
Output:
    Identifying-Backend-Bottlenecks.pptx  (project root)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Theme ----------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0xE8, 0x59, 0x0C)   # the "trace bar" orange
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
TEXT = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x5B, 0x66, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h, fill=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background()
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=8, line_spacing=1.05):
    """runs: list of paragraphs. Each paragraph = list of (text, size, bold, color, level)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        if isinstance(para, tuple):
            para = [para]
        for (txt, size, bold, color, level) in para:
            p.level = level
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = color
    return tb


def _footer(slide, idx):
    _text(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.4),
          [[("Identifying Backend Bottlenecks  ·  Brownbag", 10, False, MUTED, 0)]])
    _text(slide, Inches(11.8), Inches(7.05), Inches(1.0), Inches(0.4),
          [[(str(idx), 10, True, MUTED, 0)]], align=PP_ALIGN.RIGHT)


_page = {"n": 0}


def header(slide, title, kicker=None):
    _box(slide, 0, 0, SW, Inches(1.15), NAVY)
    _box(slide, 0, Inches(1.15), SW, Inches(0.08), ACCENT)
    if kicker:
        _text(slide, Inches(0.6), Inches(0.16), Inches(12), Inches(0.3),
              [[(kicker.upper(), 12, True, ACCENT, 0)]])
        _text(slide, Inches(0.6), Inches(0.46), Inches(12), Inches(0.7),
              [[(title, 27, True, WHITE, 0)]])
    else:
        _text(slide, Inches(0.6), Inches(0.0), Inches(12), Inches(1.15),
              [[(title, 28, True, WHITE, 0)]], anchor=MSO_ANCHOR.MIDDLE)


def content_slide(title, bullets, kicker=None, note=None):
    """bullets: list of (text, level)."""
    _page["n"] += 1
    s = prs.slides.add_slide(BLANK)
    _set_bg(s, WHITE)
    header(s, title, kicker)
    runs = []
    for txt, level in bullets:
        if level == 0:
            runs.append([("•  ", 20, True, ACCENT, 0), (txt, 20, True, TEXT, 0)])
        else:
            runs.append([("–  ", 16, False, MUTED, 1), (txt, 16, False, MUTED, 1)])
    _text(s, Inches(0.75), Inches(1.5), Inches(11.9), Inches(5.2), runs,
          space_after=9, line_spacing=1.06)
    if note:
        _box(s, Inches(0.75), Inches(6.35), Inches(11.85), Inches(0.6), LIGHT)
        _text(s, Inches(0.95), Inches(6.4), Inches(11.5), Inches(0.5),
              [[("→  ", 13, True, ACCENT, 0), (note, 13, False, TEXT, 0)]],
              anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, _page["n"])
    return s


def two_col_slide(title, left_head, left_items, right_head, right_items,
                  kicker=None, left_color=RED, right_color=GREEN):
    _page["n"] += 1
    s = prs.slides.add_slide(BLANK)
    _set_bg(s, WHITE)
    header(s, title, kicker)
    colw = Inches(5.85)
    for x, head, items, col in [
        (Inches(0.7), left_head, left_items, left_color),
        (Inches(6.85), right_head, right_items, right_color),
    ]:
        _box(s, x, Inches(1.5), colw, Inches(0.6), col)
        _text(s, x, Inches(1.5), colw, Inches(0.6),
              [[(head, 18, True, WHITE, 0)]], align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
        runs = []
        for txt, level in items:
            if level == 0:
                runs.append([("•  ", 16, True, col, 0), (txt, 16, True, TEXT, 0)])
            else:
                runs.append([("–  ", 14, False, MUTED, 1), (txt, 14, False, MUTED, 1)])
        _text(s, x + Inches(0.15), Inches(2.3), colw - Inches(0.3), Inches(4.4),
              runs, space_after=7)
    _footer(s, _page["n"])
    return s


def table_slide(title, headers, rows, kicker=None, note=None, col_widths=None):
    _page["n"] += 1
    s = prs.slides.add_slide(BLANK)
    _set_bg(s, WHITE)
    header(s, title, kicker)
    nrows, ncols = len(rows) + 1, len(headers)
    gx, gy = Inches(0.75), Inches(1.65)
    gw, gh = Inches(11.85), Inches(0.55) * nrows
    tbl = s.shapes.add_table(nrows, ncols, gx, gy, gw, gh).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = htext
        r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = WHITE; r.font.name = FONT
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            tfp = c.text_frame.paragraphs[0]
            r = tfp.add_run(); r.text = str(val)
            r.font.size = Pt(13); r.font.name = FONT; r.font.color.rgb = TEXT
            if j == 0:
                r.font.bold = True
    if note:
        _box(s, Inches(0.75), Inches(6.35), Inches(11.85), Inches(0.6), LIGHT)
        _text(s, Inches(0.95), Inches(6.4), Inches(11.5), Inches(0.5),
              [[("→  ", 13, True, ACCENT, 0), (note, 13, False, TEXT, 0)]],
              anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, _page["n"])
    return s


def divider(kicker, title, subtitle=None):
    _page["n"] += 1
    s = prs.slides.add_slide(BLANK)
    _set_bg(s, NAVY)
    _box(s, Inches(0.9), Inches(3.05), Inches(1.6), Inches(0.12), ACCENT)
    _text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.5),
          [[(kicker.upper(), 16, True, ACCENT, 0)]])
    _text(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(1.5),
          [[(title, 40, True, WHITE, 0)]])
    if subtitle:
        _text(s, Inches(0.9), Inches(4.6), Inches(11.5), Inches(1.0),
              [[(subtitle, 18, False, RGBColor(0xC7, 0xD0, 0xDA), 0)]])
    _footer(s, _page["n"])
    return s


# ===========================================================================
# TITLE
# ===========================================================================
_page["n"] += 1
s = prs.slides.add_slide(BLANK)
_set_bg(s, NAVY)
_box(s, 0, Inches(5.0), SW, Inches(0.14), ACCENT)
_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
      [[("BROWNBAG · BACKEND ENGINEERING", 16, True, ACCENT, 0)]])
_text(s, Inches(0.9), Inches(2.2), Inches(11.7), Inches(2.0),
      [[("Identifying Backend", 48, True, WHITE, 0)],
       [("Bottlenecks", 48, True, WHITE, 0)]], line_spacing=1.0)
_text(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.7),
      [[("Spotting technical AND process hurdles early — and turning them into action", 20, False,
         RGBColor(0xC7, 0xD0, 0xDA), 0)]])
_text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
      [[("Measure, don't guess.", 16, True, ACCENT, 0)]])

# ===========================================================================
# AGENDA
# ===========================================================================
content_slide("Agenda", [
    ("What a bottleneck is, and a mental model to find one", 0),
    ("Part 1 — Technical bottlenecks (with a live demo)", 0),
    ("N+1 queries · the expensive single query", 1),
    ("Measuring: query counts, timing, EXPLAIN, tracing, k6, golden signals", 1),
    ("Part 2 — Process bottlenecks", 0),
    ("Flow metrics · DORA · where delivery slows down", 1),
    ("From insight to action + cheat sheet", 0),
], kicker="Overview")

# ===========================================================================
# LEARNING OUTCOMES
# ===========================================================================
content_slide("What you'll walk away with", [
    ("A repeatable way to find the bottleneck instead of guessing", 0),
    ("The most common backend technical bottlenecks and their fixes", 0),
    ("Tools to prove it: response metrics, tracing, load tests, dashboards", 0),
    ("Awareness that process is often the bigger bottleneck than code", 0),
    ("A simple framework to turn an observation into a concrete action", 0),
], kicker="Goals")

# ===========================================================================
# WHAT IS A BOTTLENECK
# ===========================================================================
content_slide("What is a bottleneck?", [
    ("The single slowest stage that caps the whole system's throughput", 0),
    ("Like a funnel: widening any other part changes nothing", 1),
    ("Optimizing a non-bottleneck is wasted effort — it feels productive but moves no needle", 0),
    ("Two truths:", 0),
    ("It hides in small data — everything looks fine in dev, falls over at scale", 1),
    ("It moves — fix one, the next slowest stage becomes the new bottleneck", 1),
], kicker="Concept",
   note="Rule of thumb: find the slowest step first, fix it, then re-measure.")

# ===========================================================================
# TWO KINDS
# ===========================================================================
two_col_slide("Two kinds of bottlenecks",
    "TECHNICAL", [
        ("Lives in the code / system", 0),
        ("Slow queries, N+1, missing indexes", 1),
        ("No caching, blocking I/O, chatty calls", 1),
        ("Connection-pool / memory limits", 1),
        ("Found with profiling, tracing, metrics", 0),
    ],
    "PROCESS", [
        ("Lives in how we build & ship", 0),
        ("Slow code reviews, big PRs", 1),
        ("Flaky / slow CI, manual deploys", 1),
        ("Unclear requirements, handoffs, silos", 1),
        ("Found with flow & DORA metrics", 0),
    ],
    kicker="Framing", left_color=ACCENT, right_color=NAVY,
)

# ===========================================================================
# MENTAL MODEL
# ===========================================================================
content_slide("A mental model: the loop", [
    ("Observe — something feels slow, or an alert/SLO fires", 0),
    ("Measure — attach a number (latency, p95, query count, lead time)", 0),
    ("Locate — narrow to the one slowest stage (trace, profile, value-stream)", 0),
    ("Fix — make the smallest change that moves that number", 0),
    ("Verify — measure again; confirm it helped and didn't shift the problem", 0),
], kicker="How to think",
   note="Observe → Measure → Locate → Fix → Verify. The discipline is measuring before and after.")

# ===========================================================================
# GOLDEN SIGNALS
# ===========================================================================
content_slide("The golden signals (what to watch)", [
    ("Latency — how long requests take; track p95/p99, not just average", 0),
    ("Averages lie; the slow tail is where users feel pain", 1),
    ("Traffic — how much demand (requests/sec)", 0),
    ("Errors — rate of failed requests (5xx, timeouts)", 0),
    ("Saturation — how full the resources are (CPU, memory, pool, queue)", 0),
], kicker="Vocabulary",
   note="If you only learn one number today: p95 latency.")

# ===========================================================================
# PART 1 DIVIDER
# ===========================================================================
divider("Part 1", "Technical Bottlenecks",
        "Live demo: a Spring Boot REST API with slow vs. fast twins")

# ===========================================================================
# COMMON TECHNICAL
# ===========================================================================
content_slide("Common backend technical bottlenecks", [
    ("N+1 queries — one query per row instead of one query total", 0),
    ("Expensive single query — correlated subqueries, missing indexes, full scans", 0),
    ("No caching — recomputing or refetching the same data repeatedly", 0),
    ("Chatty I/O — many small network/DB calls instead of a batch", 0),
    ("Blocking work — sync calls that hog threads / the event loop", 0),
    ("No pagination — loading everything when the user sees 20 rows", 0),
    ("Resource limits — connection pool, memory, thread starvation", 0),
], kicker="Survey")

# ===========================================================================
# DEEP DIVE N+1
# ===========================================================================
table_slide("Deep dive 1 — the N+1 query problem",
    ["Endpoint", "What it does", "Cost (50 authors)"],
    [
        ["GET /api/slow/authors", "Load authors, then lazily load each one's books", "51 queries (1 + N)"],
        ["GET /api/fast/authors", "Load authors + books in one JOIN FETCH", "1 query"],
    ],
    kicker="Too many queries",
    col_widths=[3.6, 5.65, 2.6],
    note="Symptom: query count scales with rows. Fix: JOIN FETCH / eager batch / DataLoader.")

# ===========================================================================
# DEEP DIVE SLOW QUERY
# ===========================================================================
table_slide("Deep dive 2 — the expensive single query",
    ["Endpoint", "What it does", "Cost (8,000 sales)"],
    [
        ["GET /api/slow-query/sales-ranking", "Correlated subquery — re-scans table per row (O(N\u00b2))", "1 query, ~2,200 ms"],
        ["GET /api/fast-query/sales-ranking", "RANK() OVER (...) window function (O(N log N))", "1 query, ~50 ms"],
    ],
    kicker="One query, but huge",
    col_widths=[4.5, 5.15, 2.2],
    note="Same result, ~40x faster. It's ONE statement either way — only timing reveals it.")

# ===========================================================================
# COUNT VS COST
# ===========================================================================
two_col_slide("Count vs. cost — two different signals",
    "N+1 (count)", [
        ("Caught by query COUNT", 0),
        ("X-Sql-Statements: 51", 1),
        ("Many cheap statements", 1),
        ("Fix in the data-access layer", 1),
    ],
    "Slow query (cost)", [
        ("Caught by TIMING / EXPLAIN", 0),
        ("X-Sql-Statements: 1 (looks fine!)", 1),
        ("One internally huge statement", 1),
        ("Fix the query / add an index", 1),
    ],
    kicker="The key lesson", left_color=ACCENT, right_color=NAVY,
)

# ===========================================================================
# HOW TO MEASURE
# ===========================================================================
content_slide("How to measure a technical bottleneck", [
    ("Per-request signals — our demo adds headers:", 0),
    ("X-Sql-Statements (how many queries) · X-Elapsed-Ms (how long)", 1),
    ("EXPLAIN / EXPLAIN ANALYZE — read the query plan; look for full scans", 0),
    ("Slow-query log — let the database tell you which statements hurt", 0),
    ("Profiler / APM — find hot methods and slow downstream calls", 0),
    ("Distributed tracing — see the slow span inside a real request", 0),
    ("Load test (k6) — reveal what only appears under concurrency", 0),
], kicker="Tools")

# ===========================================================================
# TRACING
# ===========================================================================
content_slide("Distributed tracing — find the slowest span", [
    ("Each request becomes a trace; spans show where time goes", 0),
    ("Slow-query endpoint waterfall:", 0),
    ("validate-request ~0 ms  ·  compute-summary / map-to-dto a few ms", 1),
    ("load-ranking-slow ~2.1 s  ← one giant bar = the bottleneck", 1),
    ("N+1 endpoint: one HTTP span with ~51 child SELECT spans stacked under it", 0),
    ("Auto-instrumented HTTP + JDBC via the OpenTelemetry Java agent — no code changes", 0),
], kicker="See it, don't guess",
   note="A trace answers 'where is the time inside ONE request?'")

# ===========================================================================
# K6
# ===========================================================================
content_slide("Load testing with k6 — stress the bottleneck", [
    ("One curl shows one request; bottlenecks hurt under concurrency", 0),
    ("k6 ramps virtual users up, holds, ramps down", 0),
    ("Compare p95/p99 between a slow endpoint and its fast twin", 0),
    ("Slow endpoint's tail latency balloons; the fast one stays flat", 1),
    ("k6 run -e TARGET=slow-sales k6/compare.js", 0),
    ("Optional: stream results to Grafana Cloud k6 dashboards", 0),
], kicker="Under load",
   note="A red 'p95 < 500ms' threshold on the slow endpoint IS the lesson, not a bug.")

# ===========================================================================
# GOLDEN SIGNALS IN GRAFANA
# ===========================================================================
content_slide("Golden signals in Grafana Cloud", [
    ("Traces show one request; p95 is an aggregate that comes from METRICS", 0),
    ("Our agent already ships traces + metrics + logs", 0),
    ("Enable Application Observability to auto-build the RED dashboard:", 0),
    ("Rate (req/s) · Errors (%) · Duration (p50 / p95 / p99) per route", 1),
    ("Compare slow vs. fast endpoints side by side, in aggregate", 0),
], kicker="Aggregate view",
   note="Traces = where inside a request. App Observability / k6 = how bad p95 gets overall.")

# ===========================================================================
# THE DEMO
# ===========================================================================
content_slide("The live demo (try it yourself)", [
    ("Spring Boot 4 + H2, seeded on startup — clone & run, no external DB", 0),
    ("Four endpoints: slow/fast for N+1, slow/fast for the heavy query", 0),
    ("Compare headers:", 0),
    ("curl -s -D - -o /dev/null .../api/slow-query/sales-ranking | grep X-", 1),
    ("Run with tracing: ./scripts/run-with-tracing.sh → Grafana Cloud", 0),
    ("Full walkthrough in the repo README", 0),
], kicker="Hands-on")

# ===========================================================================
# ANY STACK
# ===========================================================================
content_slide("Same lesson, any stack", [
    ("The demo is Java/Spring, but these map everywhere:", 0),
    ("N+1 → NestJS/TypeORM relations or QueryBuilder joins; GraphQL DataLoader batching", 0),
    ("Slow query → avoid correlated subqueries & non-sargable filters; use window functions + indexes", 0),
    ("Caching → Redis/Valkey for hot reads; cache invalidation discipline", 0),
    ("Blocking → async/await, queues (SQS), background jobs off the request path", 0),
], kicker="Transfer it",
   note="Read EXPLAIN, batch your I/O, cache the hot path, and keep work off the request thread.")

# ===========================================================================
# PART 2 DIVIDER
# ===========================================================================
divider("Part 2", "Process Bottlenecks",
        "Often the code is fast — but the delivery pipeline is slow")

# ===========================================================================
# WHAT ARE PROCESS BOTTLENECKS
# ===========================================================================
content_slide("Why process bottlenecks matter", [
    ("Most lead time is spent WAITING, not coding", 0),
    ("A 2-hour change can take 5 days to reach production", 1),
    ("They're invisible in code reviews but obvious in the calendar", 0),
    ("Symptoms: work 'almost done' for days, recurring crunch, blocked tickets", 0),
    ("The fix is usually workflow, not effort or headcount", 0),
], kicker="Framing")

# ===========================================================================
# COMMON PROCESS
# ===========================================================================
content_slide("Common process bottlenecks", [
    ("Unclear requirements — rework after the fact", 0),
    ("Large PRs — slow, shallow reviews; risky merges", 0),
    ("Slow code review — PRs waiting hours/days for a first look", 0),
    ("Flaky or slow CI — long feedback loops, ignored failures", 0),
    ("Manual / risky deploys — releases batched and feared", 0),
    ("Handoffs & silos — one person/team is a single point of contact", 0),
    ("Too much WIP — everything in progress, nothing finished", 0),
], kicker="Survey")

# ===========================================================================
# HOW TO SPOT PROCESS
# ===========================================================================
content_slide("How to spot process bottlenecks", [
    ("Make the invisible visible with flow metrics:", 0),
    ("Lead time — idea/commit → production", 1),
    ("Cycle time — work started → done", 1),
    ("PR review time — opened → first review → merged", 1),
    ("CI duration & flake rate — pipeline speed and trust", 1),
    ("WIP — how many things are in flight at once", 1),
    ("Look for the stage where work piles up and waits", 0),
], kicker="Measure flow",
   note="Same loop as code: Observe → Measure → Locate → Fix → Verify.")

# ===========================================================================
# DORA
# ===========================================================================
table_slide("DORA — four keys of delivery performance",
    ["Metric", "What it measures", "Improve by"],
    [
        ["Deployment Frequency", "How often you ship", "Smaller batches, automation"],
        ["Lead Time for Changes", "Commit → production", "Trunk-based, fast CI"],
        ["Change Failure Rate", "% of deploys causing issues", "Tests, smaller changes"],
        ["Time to Restore", "How fast you recover", "Monitoring, rollbacks"],
    ],
    kicker="Industry signals",
    col_widths=[3.4, 5.05, 3.4],
    note="Speed and stability rise together — they are not a trade-off.")

# ===========================================================================
# INSIGHT TO ACTION
# ===========================================================================
content_slide("From insight to action", [
    ("Identify — name the one slowest stage (technical or process)", 0),
    ("Quantify — attach a number and a baseline", 0),
    ("Prioritize — rank by impact vs. effort; pick the biggest cheap win", 0),
    ("Act — make one focused change (a code fix or a workflow change)", 0),
    ("Verify — measure again; keep it if the number moved, revert if not", 0),
], kicker="The framework",
   note="One change at a time, always measured. Write the before/after number in the PR.")

# ===========================================================================
# CHEAT SHEET
# ===========================================================================
two_col_slide("Cheat sheet",
    "TECHNICAL", [
        ("Check query COUNT and TIMING", 0),
        ("Read EXPLAIN before optimizing", 0),
        ("Batch I/O; avoid N+1", 0),
        ("Cache the hot path", 0),
        ("Trace to find the slow span", 0),
        ("Load test for p95 under load", 0),
    ],
    "PROCESS", [
        ("Keep PRs small", 0),
        ("Review fast; protect review time", 0),
        ("Make CI fast and trustworthy", 0),
        ("Automate deploys", 0),
        ("Limit WIP; finish before starting", 0),
        ("Track lead time & DORA", 0),
    ],
    kicker="Keep this handy", left_color=ACCENT, right_color=NAVY,
)

# ===========================================================================
# TAKEAWAYS
# ===========================================================================
content_slide("Key takeaways", [
    ("Measure, don't guess — turn 'feels slow' into a number", 0),
    ("Fix the actual bottleneck; optimizing elsewhere is wasted effort", 0),
    ("Count and cost are different signals — check both", 0),
    ("The biggest win is often process, not code", 0),
    ("Make it a habit: Observe → Measure → Locate → Fix → Verify", 0),
], kicker="Remember")

# ===========================================================================
# THANK YOU
# ===========================================================================
_page["n"] += 1
s = prs.slides.add_slide(BLANK)
_set_bg(s, NAVY)
_box(s, 0, Inches(4.7), SW, Inches(0.14), ACCENT)
_text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.2),
      [[("Thank you — questions?", 40, True, WHITE, 0)]])
_text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.4),
      [[("Demo repo: bottleneck-springboot (README has full setup)", 18, False, RGBColor(0xC7,0xD0,0xDA), 0)],
       [("Clone it, run the slow vs. fast endpoints, and watch the numbers.", 16, False, RGBColor(0xC7,0xD0,0xDA), 0)]])

out = "Identifying-Backend-Bottlenecks.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
